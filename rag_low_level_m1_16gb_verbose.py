"""
rag_low_level_m1_16gb_verbose.py

Goal
-----
A low-level RAG pipeline, fully local, with Postgres+pgvector for vectors, and llama.cpp GGUF for generation.
This is tuned to be "reasonable" on a 16GB Mac mini M1.

What you will learn by reading logs
-----------------------------------
1) How many Documents the PDF became (often pages)
2) How many chunks were produced + why overlap matters
3) How embeddings are computed (batched) + time per batch
4) How many rows are stored in Postgres + table reset behavior
5) What retrieval returns (scores, metadata, text previews)
6) What the LLM answers given retrieved evidence
"""

import os
import sys
import time
import platform
import logging
import argparse
import hashlib
import json
import re
from datetime import datetime
from dataclasses import dataclass, field
from typing import Any, List, Optional, Iterable, Tuple
from pathlib import Path

import psycopg2
from psycopg2 import OperationalError as PgOperationalError

try:
    from tqdm import tqdm
except ImportError:
    tqdm = None
    logging.warning("tqdm not installed, progress bars disabled. Install with: pip install tqdm")

from llama_index.readers.file import PyMuPDFReader
from llama_index.core.node_parser import SentenceSplitter
from llama_index.core.schema import TextNode, NodeWithScore
from llama_index.embeddings.huggingface import HuggingFaceEmbedding

from llama_index.vector_stores.postgres import PGVectorStore
from llama_index.core.vector_stores import VectorStoreQuery
from llama_index.core import QueryBundle
from llama_index.core.retrievers import BaseRetriever
from llama_index.core.query_engine import RetrieverQueryEngine
from llama_index.llms.llama_cpp import LlamaCPP

# Optional: vLLM for GPU-accelerated inference (15x faster than llama.cpp CPU)
try:
    from vllm_wrapper import build_vllm_llm
    VLLM_AVAILABLE = True
except ImportError:
    VLLM_AVAILABLE = False

# vLLM Server Mode (best performance - no reload between queries)
try:
    from vllm_client import build_vllm_client
    VLLM_CLIENT_AVAILABLE = True
except ImportError:
    VLLM_CLIENT_AVAILABLE = False


# -----------------------
# Optional memory logging (nice to have, not required)
# -----------------------
try:
    import psutil  # pip install psutil
except ImportError as e:
    psutil = None
    # Note: psutil import failed, memory stats will not be available
except Exception as e:
    psutil = None
    logging.warning(f"Unexpected error importing psutil: {type(e).__name__}: {e}")


# -----------------------
# Logging configuration
# -----------------------
def setup_logging() -> logging.Logger:
    """
    LOG_LEVEL can be: DEBUG, INFO, WARNING, ERROR
    Default INFO is already pretty chatty.
    """
    level_str = os.getenv("LOG_LEVEL", "INFO").upper()
    level = getattr(logging, level_str, logging.INFO)

    logging.basicConfig(
        level=level,
        format="%(asctime)s | %(levelname)-7s | %(message)s",
        datefmt="%H:%M:%S",
    )
    return logging.getLogger("rag")


log = setup_logging()


def log_system_info():
    log.info("System info:")
    log.info(f"  Python: {sys.version.split()[0]}")
    log.info(f"  Platform: {platform.platform()}")
    if psutil:
        vm = psutil.virtual_memory()
        log.info(f"  RAM: total={vm.total/1e9:.1f}GB available={vm.available/1e9:.1f}GB used={vm.percent}%")
    else:
        log.info("  RAM: (psutil not installed; skipping memory stats)")


def now_ms() -> int:
    return int(time.time() * 1000)


def dur_s(start_ms: int) -> float:
    return (now_ms() - start_ms) / 1000.0


def chunked(it: List[Any], n: int) -> Iterable[List[Any]]:
    """Yield list chunks of size n."""
    for i in range(0, len(it), n):
        yield it[i : i + n]


def preview(text: str, n: int = 220) -> str:
    """Small helper to keep logs readable.
    Set LOG_FULL_CHUNKS=1 to disable truncation."""
    # Check if full chunks should be displayed
    if os.getenv("LOG_FULL_CHUNKS", "0") == "1":
        return (text or "").strip()

    t = (text or "").replace("\n", " ").strip()
    return (t[:n] + "…") if len(t) > n else t


def colorize_participants(text: str) -> str:
    """
    Colorize participant names in chat logs for better readability.
    Detects patterns like [date time] Name: message and assigns colors.
    Set COLORIZE_CHUNKS=1 to enable.
    """
    if os.getenv("COLORIZE_CHUNKS", "0") != "1":
        return text

    import re

    # ANSI color codes
    COLORS = [
        '\033[91m',  # Red
        '\033[92m',  # Green
        '\033[93m',  # Yellow
        '\033[94m',  # Blue
        '\033[95m',  # Magenta
        '\033[96m',  # Cyan
        '\033[97m',  # White
        '\033[31m',  # Dark Red
        '\033[32m',  # Dark Green
        '\033[33m',  # Dark Yellow
        '\033[34m',  # Dark Blue
        '\033[35m',  # Dark Magenta
        '\033[36m',  # Dark Cyan
    ]
    RESET = '\033[0m'
    BOLD = '\033[1m'

    # Find all participant names in pattern: [date time] Name:
    # Match patterns like [2022-10-13 13:22] EB: or [2024-06-18 23:40] Arnaud Grd:
    pattern = r'\[[\d\-]+ [\d:]+\]\s+([^:]+):'
    matches = re.findall(pattern, text)

    if not matches:
        return text

    # Get unique participants and assign colors
    participants = list(set(matches))
    participant_colors = {}
    for i, participant in enumerate(sorted(participants)):
        participant_colors[participant] = COLORS[i % len(COLORS)]

    # Colorize each participant name in the text
    colored_text = text
    for participant, color in participant_colors.items():
        # Use word boundaries to avoid partial matches
        # Replace "Name:" with colored version
        colored_text = re.sub(
            rf'(\[[\d\-]+ [\d:]+\]\s+)({re.escape(participant)})(:)',
            rf'\1{BOLD}{color}\2{RESET}\3',
            colored_text
        )

    return colored_text


def compute_file_hash(file_path: str) -> str:
    """Compute SHA-256 hash of a file for tracking."""
    sha256 = hashlib.sha256()
    with open(file_path, "rb") as f:
        for chunk in iter(lambda: f.read(8192), b""):
            sha256.update(chunk)
    return sha256.hexdigest()


def extract_chat_metadata(text: str) -> dict:
    """
    Extract metadata from chat log messages.
    Detects patterns like: [2022-10-13 13:22] Name: message

    Returns dict with:
        - participants: List of unique participant names
        - dates: List of dates mentioned
        - date_range: Tuple of (earliest, latest) dates
        - message_count: Number of messages in chunk
    """
    import re
    from datetime import datetime

    # Pattern: [YYYY-MM-DD HH:MM] Name: message
    pattern = r'\[(\d{4}-\d{2}-\d{2})\s+(\d{2}:\d{2})\]\s+([^:]+):'
    matches = re.findall(pattern, text)

    if not matches:
        return {}

    participants = []
    dates = []

    for date_str, time_str, participant in matches:
        participants.append(participant.strip())
        dates.append(date_str)

    # Get unique values
    unique_participants = list(set(participants))
    unique_dates = sorted(set(dates))

    metadata = {
        "participants": unique_participants,
        "participant_count": len(unique_participants),
        "message_count": len(matches),
        "is_chat_log": True,
    }

    # Add date range if dates found
    if unique_dates:
        metadata["dates"] = unique_dates
        metadata["earliest_date"] = unique_dates[0]
        metadata["latest_date"] = unique_dates[-1]
        metadata["date_range"] = f"{unique_dates[0]} to {unique_dates[-1]}"

    # Add dominant participant (most messages in chunk)
    if participants:
        from collections import Counter
        participant_counts = Counter(participants)
        dominant = participant_counts.most_common(1)[0]
        metadata["dominant_participant"] = dominant[0]
        metadata["dominant_participant_count"] = dominant[1]

    return metadata


def clean_html_content(html: str) -> str:
    """
    Clean HTML content for better embedding quality.
    Extracts text, removes scripts/styles, normalizes whitespace.
    """
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        # Fallback: basic regex cleaning if BeautifulSoup not available
        import re
        # Remove script and style blocks
        text = re.sub(r'<script[^>]*>.*?</script>', '', html, flags=re.DOTALL | re.IGNORECASE)
        text = re.sub(r'<style[^>]*>.*?</style>', '', text, flags=re.DOTALL | re.IGNORECASE)
        # Remove HTML tags
        text = re.sub(r'<[^>]+>', ' ', text)
        # Decode HTML entities
        import html as html_module
        text = html_module.unescape(text)
        # Normalize whitespace
        text = ' '.join(text.split())
        return text

    # Use BeautifulSoup for better parsing
    soup = BeautifulSoup(html, 'html.parser')

    # Remove unwanted elements
    for element in soup(['script', 'style', 'head', 'meta', 'link', 'noscript']):
        element.decompose()

    # Extract text with newlines preserved for structure
    text = soup.get_text(separator='\n')

    # Clean up whitespace while preserving paragraph structure
    lines = []
    for line in text.splitlines():
        line = line.strip()
        if line:  # Skip empty lines
            # Collapse multiple spaces within a line
            line = ' '.join(line.split())
            lines.append(line)

    # Join with single newlines, collapse multiple blank lines
    text = '\n'.join(lines)

    return text


def parse_args() -> argparse.Namespace:
    """Parse command-line arguments."""
    parser = argparse.ArgumentParser(
        description="Local RAG pipeline with PostgreSQL + pgvector",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  # Full pipeline (index + query)
  python %(prog)s

  # Query existing index only (skip re-indexing)
  python %(prog)s --query-only

  # Interactive mode (ask multiple questions)
  python %(prog)s --interactive

  # Single query via CLI
  python %(prog)s --query "What are the main findings?"

  # Use different document
  PDF_PATH=my_doc.pdf PGTABLE=my_doc python %(prog)s

Environment Variables:
  See README.md for full list of configuration options
        """
    )

    parser.add_argument(
        "--query-only",
        action="store_true",
        help="Skip document ingestion, only run query on existing index"
    )

    parser.add_argument(
        "--index-only",
        action="store_true",
        help="Only index documents, skip querying (useful for batch indexing)"
    )

    parser.add_argument(
        "--interactive",
        "-i",
        action="store_true",
        help="Interactive REPL mode for multiple queries"
    )

    parser.add_argument(
        "--query",
        "-q",
        type=str,
        help="Single query to run (overrides QUESTION env var)"
    )

    parser.add_argument(
        "--doc",
        "-d",
        type=str,
        help="Document path (overrides PDF_PATH env var)"
    )

    parser.add_argument(
        "--skip-validation",
        action="store_true",
        help="Skip startup validation (use with caution)"
    )

    return parser.parse_args()


# -----------------------
# Table Name Generation
# -----------------------
def sanitize_name(name: str, max_length: int = 30) -> str:
    """
    Sanitize a filename/folder for use in table names.

    Args:
        name: Original filename or folder name
        max_length: Maximum length for sanitized name

    Returns:
        Sanitized name safe for PostgreSQL table names
    """
    # Remove file extension
    name = Path(name).stem

    # Replace problematic characters with underscores
    name = name.replace('-', '_').replace(' ', '_').replace('.', '_')

    # Remove any remaining non-alphanumeric characters except underscores
    name = re.sub(r'[^a-zA-Z0-9_]', '', name)

    # Ensure it starts with a letter (PostgreSQL requirement)
    if name and not name[0].isalpha():
        name = 'doc_' + name

    # Truncate to max length
    if len(name) > max_length:
        name = name[:max_length]

    return name.lower()


def extract_model_short_name(model_name: str) -> str:
    """
    Extract a short, readable name from embedding model path.

    Examples:
        "BAAI/bge-small-en" -> "bge"
        "sentence-transformers/all-MiniLM-L6-v2" -> "minilm"
        "intfloat/multilingual-e5-small" -> "e5"

    Args:
        model_name: Full model name/path

    Returns:
        Short model identifier
    """
    # Extract the last part after /
    name = model_name.split('/')[-1]

    # Common patterns
    if 'bge' in name.lower():
        return 'bge'
    elif 'minilm' in name.lower():
        return 'minilm'
    elif 'e5' in name.lower():
        return 'e5'
    elif 'mpnet' in name.lower():
        return 'mpnet'
    elif 'roberta' in name.lower():
        return 'roberta'
    elif 'bert' in name.lower():
        return 'bert'
    else:
        # Take first meaningful word (remove common prefixes)
        parts = name.lower().replace('sentence-', '').replace('all-', '').split('-')
        return parts[0][:8]


def generate_table_name(
    pdf_path: str,
    chunk_size: int,
    chunk_overlap: int,
    embed_model: str
) -> str:
    """
    Generate a descriptive table name from configuration.

    Format: {file/folder}_cs{chunk_size}_ov{overlap}_{model}_{YYMMDD}

    Args:
        pdf_path: Path to document or folder
        chunk_size: Chunk size parameter
        chunk_overlap: Chunk overlap parameter
        embed_model: Embedding model name

    Returns:
        Generated table name

    Examples:
        generate_table_name("data/inbox_clean", 700, 150, "BAAI/bge-small-en")
        -> "inbox_clean_cs700_ov150_bge_251219"

        generate_table_name("data/report.pdf", 500, 100, "intfloat/e5-small")
        -> "report_cs500_ov100_e5_251219"
    """
    # Extract and sanitize document name
    path = Path(pdf_path)
    if path.is_dir():
        doc_name = sanitize_name(path.name)
    else:
        doc_name = sanitize_name(path.stem)

    # Extract short model name
    model_short = extract_model_short_name(embed_model)

    # Get date in YYMMDD format
    date_str = datetime.now().strftime("%y%m%d")

    # Construct table name
    table_name = f"{doc_name}_cs{chunk_size}_ov{chunk_overlap}_{model_short}_{date_str}"

    return table_name


# -----------------------
# Configuration (all overrideable via env vars)
# -----------------------
@dataclass
class Settings:
    # Postgres
    db_name: str = os.getenv("DB_NAME", "vector_db")
    host: str = os.getenv("PGHOST", "localhost")
    port: str = os.getenv("PGPORT", "5432")
    user: str = os.getenv("PGUSER")
    password: str = os.getenv("PGPASSWORD")
    table: str = ""  # Will be auto-generated if not set via PGTABLE

    # Input
    pdf_path: str = os.getenv("PDF_PATH", "data/llama2.pdf")

    # Reset behaviors
    # RESET_TABLE=1 is useful while iterating so you don't duplicate rows every run
    reset_table: bool = os.getenv("RESET_TABLE", "0") == "1"
    # RESET_DB=1 is more nuclear; only use if you want a fresh DB
    reset_db: bool = os.getenv("RESET_DB", "0") == "1"

    # Chunking knobs (RAG quality knobs)
    chunk_size: int = int(os.getenv("CHUNK_SIZE", "700"))
    chunk_overlap: int = int(os.getenv("CHUNK_OVERLAP", "150"))

    # Retrieval knobs
    top_k: int = int(os.getenv("TOP_K", "4"))

    # Advanced retrieval features
    # HYBRID_ALPHA: 0.0=pure BM25, 0.5=balanced, 1.0=pure vector (default)
    hybrid_alpha: float = float(os.getenv("HYBRID_ALPHA", "1.0"))
    # ENABLE_FILTERS: Enable metadata-based filtering (participant, date)
    enable_filters: bool = os.getenv("ENABLE_FILTERS", "1") == "1"
    # MMR_THRESHOLD: Enable MMR diversity (0=disabled, 0.5=balanced, 1.0=max relevance)
    mmr_threshold: float = float(os.getenv("MMR_THRESHOLD", "0.0"))

    # Embeddings knobs
    embed_model_name: str = os.getenv("EMBED_MODEL", "BAAI/bge-small-en")
    embed_dim: int = int(os.getenv("EMBED_DIM", "384"))
    # Increased default batch size for better throughput (MLX: 64-128, HuggingFace: 32)
    embed_batch: int = int(os.getenv("EMBED_BATCH", "32"))  # Was: 16
    # EMBED_BACKEND: huggingface (default) | mlx (Apple Silicon optimized, 5-20x faster)
    embed_backend: str = os.getenv("EMBED_BACKEND", "huggingface")

    # LLM knobs (llama.cpp)
    # Default: Mistral 7B Instruct GGUF Q4_K_M (good for 16GB M1)
    model_url: str = os.getenv(
        "MODEL_URL",
        "https://huggingface.co/TheBloke/Mistral-7B-Instruct-v0.2-GGUF/resolve/main/"
        "mistral-7b-instruct-v0.2.Q4_K_M.gguf",
    )
    # If you download manually, set MODEL_PATH to a local file and it will skip model_url.
    model_path: str = os.getenv("MODEL_PATH", "")

    temperature: float = float(os.getenv("TEMP", "0.1"))
    max_new_tokens: int = int(os.getenv("MAX_NEW_TOKENS", "256"))
    context_window: int = int(os.getenv("CTX", "3072"))

    # On Apple Silicon, these matter a lot:
    # - N_GPU_LAYERS: offload more layers to Metal can speed up, but too high can crash or thrash
    # - N_BATCH: affects prompt processing throughput + peak memory
    n_gpu_layers: int = int(os.getenv("N_GPU_LAYERS", "16"))
    n_batch: int = int(os.getenv("N_BATCH", "128"))

    # Question
    question: str = os.getenv(
        "QUESTION",
        "Summarize the key safety-related training ideas described in this paper.",
    )

    def __post_init__(self):
        """Validate credentials are set after dataclass initialization."""
        if not self.user or not self.password:
            raise ValueError(
                "Database credentials not set!\n"
                "Required environment variables:\n"
                "  PGUSER=your_database_user\n"
                "  PGPASSWORD=your_database_password\n"
                "Set them in .env file or export them."
            )

    def validate(self) -> None:
        """
        Validate settings and provide helpful error messages.
        Raises ValueError with actionable error messages if validation fails.
        """
        errors = []

        # Validate document path (file or folder)
        doc_path = Path(self.pdf_path)
        supported_extensions = {
            ".pdf", ".docx", ".pptx",
            ".txt", ".md", ".html", ".htm", ".json", ".csv", ".xml", ".xsl",
            ".py", ".js", ".ts", ".jsx", ".tsx", ".cpp", ".c", ".h", ".hpp",
            ".java", ".rb", ".go", ".rs", ".php", ".m", ".swift", ".kt",
            ".sh", ".bash", ".zsh", ".yaml", ".yml", ".toml", ".ini", ".cfg",
            ".sql", ".r", ".scala", ".pl", ".lua", ".ex", ".exs", ".clj",
        }
        if not doc_path.exists():
            errors.append(
                f"Document path not found: {self.pdf_path}\n"
                f"  Fix: Set PDF_PATH environment variable to an existing file or folder"
            )
        elif doc_path.is_file() and doc_path.suffix.lower() not in supported_extensions:
            errors.append(
                f"Unsupported file format: {self.pdf_path}\n"
                f"  Extension: {doc_path.suffix}\n"
                f"  Supported: .pdf, .docx, .pptx, .html, .json, .txt, .md, .py, .cpp, etc."
            )

        # Validate chunk settings
        if self.chunk_size <= 0:
            errors.append(f"CHUNK_SIZE must be positive, got: {self.chunk_size}")
        if self.chunk_overlap < 0:
            errors.append(f"CHUNK_OVERLAP cannot be negative, got: {self.chunk_overlap}")
        if self.chunk_overlap >= self.chunk_size:
            errors.append(
                f"CHUNK_OVERLAP ({self.chunk_overlap}) must be less than CHUNK_SIZE ({self.chunk_size})\n"
                f"  Recommended: CHUNK_OVERLAP should be 10-20% of CHUNK_SIZE"
            )

        # Validate retrieval settings
        if self.top_k <= 0:
            errors.append(f"TOP_K must be positive, got: {self.top_k}")

        # Validate embedding settings
        if self.embed_dim <= 0:
            errors.append(f"EMBED_DIM must be positive, got: {self.embed_dim}")
        if self.embed_batch <= 0:
            errors.append(f"EMBED_BATCH must be positive, got: {self.embed_batch}")

        # Validate LLM settings
        if self.temperature < 0 or self.temperature > 2:
            errors.append(
                f"TEMP should be between 0 and 2, got: {self.temperature}\n"
                f"  Recommended: 0.0-0.1 for factual RAG, 0.7-1.0 for creative tasks"
            )
        if self.max_new_tokens <= 0:
            errors.append(f"MAX_NEW_TOKENS must be positive, got: {self.max_new_tokens}")
        if self.context_window <= 0:
            errors.append(f"CTX must be positive, got: {self.context_window}")
        if self.n_gpu_layers < 0:
            errors.append(f"N_GPU_LAYERS cannot be negative, got: {self.n_gpu_layers}")
        if self.n_batch <= 0:
            errors.append(f"N_BATCH must be positive, got: {self.n_batch}")

        # Validate model configuration
        if self.model_path and not Path(self.model_path).exists():
            errors.append(
                f"MODEL_PATH specified but file not found: {self.model_path}\n"
                f"  Fix: Either download the model or unset MODEL_PATH to auto-download"
            )

        # Validate database settings
        if not self.db_name or not self.db_name.strip():
            errors.append("DB_NAME cannot be empty")
        if not self.table or not self.table.strip():
            errors.append("PGTABLE cannot be empty")

        try:
            port_int = int(self.port)
            if port_int < 1 or port_int > 65535:
                errors.append(f"PGPORT must be between 1 and 65535, got: {port_int}")
        except ValueError:
            errors.append(f"PGPORT must be a valid integer, got: {self.port}")

        if errors:
            error_msg = "\n\n".join([f"❌ {err}" for err in errors])
            raise ValueError(
                f"\n\n{'='*70}\n"
                f"Configuration Validation Failed\n"
                f"{'='*70}\n\n"
                f"{error_msg}\n\n"
                f"{'='*70}\n"
            )


S = Settings()

# Auto-generate table name if PGTABLE was not explicitly set
if not os.getenv("PGTABLE"):
    S.table = generate_table_name(
        S.pdf_path,
        S.chunk_size,
        S.chunk_overlap,
        S.embed_model_name
    )
    log.debug(f"Auto-generated table name: {S.table}")
else:
    S.table = os.getenv("PGTABLE")

# Validation will be called in main() after CLI args are parsed


# -----------------------
# DB helpers with retry logic
# -----------------------
def retry_with_backoff(func, max_retries=3, initial_delay=1.0, backoff_factor=2.0, exception_types=(Exception,)):
    """
    Retry a function with exponential backoff.

    Args:
        func: Function to retry (should take no arguments)
        max_retries: Maximum number of retry attempts
        initial_delay: Initial delay in seconds
        backoff_factor: Multiplier for delay after each retry
        exception_types: Tuple of exception types to catch and retry

    Returns:
        Result of successful function call

    Raises:
        Last exception if all retries fail
    """
    delay = initial_delay
    last_exception = None

    for attempt in range(max_retries):
        try:
            return func()
        except exception_types as e:
            last_exception = e
            if attempt < max_retries - 1:
                log.warning(
                    f"Attempt {attempt + 1}/{max_retries} failed: {type(e).__name__}: {e}. "
                    f"Retrying in {delay:.1f}s..."
                )
                time.sleep(delay)
                delay *= backoff_factor
            else:
                log.error(f"All {max_retries} attempts failed. Giving up.")

    raise last_exception


def admin_conn():
    """
    Connect to the 'postgres' admin DB (useful to create/drop databases).
    Includes retry logic with helpful error messages.
    """
    def _connect():
        try:
            return psycopg2.connect(
                dbname="postgres",
                host=S.host,
                port=S.port,
                user=S.user,
                password=S.password,
                connect_timeout=10,
            )
        except PgOperationalError as e:
            error_msg = str(e).lower()
            if "could not connect" in error_msg or "connection refused" in error_msg:
                raise PgOperationalError(
                    f"Cannot connect to PostgreSQL at {S.host}:{S.port}\n"
                    f"  Fix: Make sure PostgreSQL is running:\n"
                    f"    docker-compose up -d\n"
                    f"    docker-compose ps  # Check status\n"
                    f"  Original error: {e}"
                )
            elif "authentication failed" in error_msg or "password" in error_msg:
                raise PgOperationalError(
                    f"PostgreSQL authentication failed for user '{S.user}'\n"
                    f"  Fix: Check PGUSER and PGPASSWORD environment variables\n"
                    f"  Original error: {e}"
                )
            else:
                raise

    return retry_with_backoff(_connect, max_retries=3, exception_types=(PgOperationalError,))


def db_conn():
    """
    Connect to the target DB where pgvector tables live.
    Includes retry logic with helpful error messages.
    """
    def _connect():
        try:
            return psycopg2.connect(
                dbname=S.db_name,
                host=S.host,
                port=S.port,
                user=S.user,
                password=S.password,
                connect_timeout=10,
            )
        except PgOperationalError as e:
            error_msg = str(e).lower()
            if "database" in error_msg and "does not exist" in error_msg:
                raise PgOperationalError(
                    f"Database '{S.db_name}' does not exist\n"
                    f"  Fix: The script will try to create it automatically, or create manually:\n"
                    f"    docker-compose exec db psql -U {S.user} -c 'CREATE DATABASE {S.db_name};'\n"
                    f"  Original error: {e}"
                )
            elif "could not connect" in error_msg or "connection refused" in error_msg:
                raise PgOperationalError(
                    f"Cannot connect to PostgreSQL at {S.host}:{S.port}\n"
                    f"  Fix: Make sure PostgreSQL is running:\n"
                    f"    docker-compose up -d\n"
                    f"  Original error: {e}"
                )
            else:
                raise

    return retry_with_backoff(_connect, max_retries=3, exception_types=(PgOperationalError,))


def ensure_db_exists():
    """
    For Docker Compose setups where POSTGRES_DB already creates the DB, this is unnecessary.
    But it's safe-ish if user has permissions.
    """
    start = now_ms()
    try:
        conn = admin_conn()
        conn.autocommit = True
        with conn.cursor() as c:
            if S.reset_db:
                log.warning(f"RESET_DB=1 -> Dropping database {S.db_name} (data loss).")
                c.execute(f"DROP DATABASE IF EXISTS {S.db_name}")
            # Create DB (will fail if exists -> we catch below)
            c.execute(f"CREATE DATABASE {S.db_name}")
        conn.close()
        log.info(f"DB ensure/create done in {dur_s(start):.2f}s (created new DB).")
    except Exception as e:
        # Most common: "already exists" or no permission.
        log.info(f"DB ensure/create skipped/failed harmlessly: {type(e).__name__}: {e} ({dur_s(start):.2f}s)")


def ensure_pgvector_extension():
    """
    pgvector must be enabled per database:
      CREATE EXTENSION vector;
    Without it, type 'vector' doesn't exist (your earlier error).
    """
    start = now_ms()
    conn = db_conn()
    conn.autocommit = True
    with conn.cursor() as c:
        c.execute("CREATE EXTENSION IF NOT EXISTS vector;")
    conn.close()
    log.info(f"pgvector extension ensured in {dur_s(start):.2f}s")


def reset_table_if_requested():
    """
    RESET_TABLE=1 drops the vector store table so re-running does not duplicate rows.
    """
    if not S.reset_table:
        log.info("RESET_TABLE=0 -> keeping existing table (may duplicate on re-ingest).")
        return

    start = now_ms()
    conn = db_conn()
    conn.autocommit = True
    with conn.cursor() as c:
        log.warning(f"RESET_TABLE=1 -> Dropping table '{S.table}' if it exists.")
        c.execute(f'DROP TABLE IF EXISTS "{S.table}";')
    conn.close()
    log.info(f"Table reset done in {dur_s(start):.2f}s")


def count_rows() -> Optional[int]:
    """
    Useful to see ingestion effect. If table doesn't exist yet, return None.
    Note: PGVectorStore adds 'data_' prefix to table names.
    """
    try:
        conn = db_conn()
        # PGVectorStore uses "data_{table_name}" format
        actual_table = f"data_{S.table}"
        with conn.cursor() as c:
            c.execute(f'SELECT COUNT(*) FROM "{actual_table}";')
            n = int(c.fetchone()[0])
        conn.close()
        return n
    except psycopg2.errors.UndefinedTable:
        # Table doesn't exist yet - this is expected on first run
        log.debug(f"Table 'data_{S.table}' does not exist yet (normal for first run)")
        return None
    except PgOperationalError as e:
        log.warning(f"Failed to count rows (database connection issue): {e}")
        return None
    except Exception as e:
        log.warning(f"Failed to count rows: {type(e).__name__}: {e}")
        return None


def check_index_configuration() -> Optional[dict]:
    """
    Check the configuration of existing index by sampling stored metadata.
    Returns dict with configuration info, or None if table doesn't exist or is empty.

    This helps detect "mixed index" scenarios where the table contains chunks
    created with different chunk_size/overlap settings.
    """
    try:
        conn = db_conn()
        actual_table = f"data_{S.table}"

        with conn.cursor() as c:
            # Sample some rows to check their metadata
            # We check both old and new rows in case there's a mix
            c.execute(f'''
                SELECT metadata, id FROM "{actual_table}"
                ORDER BY id LIMIT 10
            ''')
            rows = c.fetchall()

        conn.close()

        if not rows:
            return None

        # Extract configurations from sampled rows
        configs = []
        for metadata, row_id in rows:
            if metadata and isinstance(metadata, dict):
                config = {
                    "chunk_size": metadata.get("_chunk_size"),
                    "chunk_overlap": metadata.get("_chunk_overlap"),
                    "embed_model": metadata.get("_embed_model"),
                    "index_signature": metadata.get("_index_signature"),
                }
                configs.append(config)

        if not configs:
            # Old index without metadata - can't determine configuration
            return {"legacy": True, "has_metadata": False}

        # Check if all configs are the same
        first_config = configs[0]
        all_same = all(c == first_config for c in configs)

        result = {
            "legacy": False,
            "has_metadata": True,
            "chunk_size": first_config.get("chunk_size"),
            "chunk_overlap": first_config.get("chunk_overlap"),
            "embed_model": first_config.get("embed_model"),
            "index_signature": first_config.get("index_signature"),
            "is_consistent": all_same,
            "sampled_configs": len(set(str(c) for c in configs)),
        }

        return result

    except psycopg2.errors.UndefinedTable:
        return None
    except Exception as e:
        log.warning(f"Failed to check index configuration: {type(e).__name__}: {e}")
        return None


def save_query_log(
    question: str,
    answer: str,
    retrieved_chunks: List[Any],
    retrieval_time: float,
    generation_time: float,
    parameters: dict
) -> Optional[str]:
    """
    Save query results to a structured JSON log file.

    Args:
        question: The query text
        answer: The generated answer
        retrieved_chunks: List of retrieved NodeWithScore objects
        retrieval_time: Time taken for retrieval (seconds)
        generation_time: Time taken for generation (seconds)
        parameters: Dictionary of all relevant parameters

    Returns:
        Path to the log file, or None if logging is disabled
    """
    # Check if logging is enabled
    if os.getenv("LOG_QUERIES", "0") != "1":
        return None

    try:
        # Create log directory structure
        log_dir = Path("query_logs") / S.table
        log_dir.mkdir(parents=True, exist_ok=True)

        # Generate timestamp-based filename
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S_%f")
        log_file = log_dir / f"{timestamp}_query.json"

        # Process retrieved chunks
        chunks_data = []
        for i, chunk in enumerate(retrieved_chunks):
            chunk_info = {
                "rank": i + 1,
                "similarity_score": float(chunk.score) if hasattr(chunk, 'score') and chunk.score is not None else None,
                "text": chunk.node.get_content(),
                "text_length": len(chunk.node.get_content()),
                "metadata": chunk.node.metadata if hasattr(chunk.node, 'metadata') else {}
            }
            chunks_data.append(chunk_info)

        # Calculate retrieval quality metrics
        scores = [c["similarity_score"] for c in chunks_data if c["similarity_score"] is not None]
        quality_metrics = {}
        if scores:
            quality_metrics = {
                "best_score": max(scores),
                "worst_score": min(scores),
                "average_score": sum(scores) / len(scores),
                "score_range": max(scores) - min(scores)
            }

        # Build complete log structure
        log_data = {
            "metadata": {
                "timestamp": datetime.now().isoformat(),
                "document_table": S.table,
                "pdf_path": S.pdf_path,
                "log_file": str(log_file)
            },
            "query": {
                "question": question,
                "question_length": len(question),
                "question_words": len(question.split())
            },
            "parameters": parameters,
            "retrieval": {
                "retrieved_chunks": chunks_data,
                "num_chunks": len(chunks_data),
                "retrieval_time_seconds": round(retrieval_time, 3),
                "quality_metrics": quality_metrics
            },
            "generation": {
                "answer": answer,
                "answer_length": len(answer),
                "answer_words": len(answer.split()),
                "generation_time_seconds": round(generation_time, 3),
                "tokens_per_second": round(len(answer.split()) / generation_time, 2) if generation_time > 0 else 0
            },
            "performance": {
                "total_time_seconds": round(retrieval_time + generation_time, 3),
                "retrieval_percentage": round((retrieval_time / (retrieval_time + generation_time)) * 100, 1) if (retrieval_time + generation_time) > 0 else 0,
                "generation_percentage": round((generation_time / (retrieval_time + generation_time)) * 100, 1) if (retrieval_time + generation_time) > 0 else 0
            }
        }

        # Write to file with pretty formatting
        with open(log_file, 'w', encoding='utf-8') as f:
            json.dump(log_data, f, indent=2, ensure_ascii=False)

        log.info(f"\n💾 Query log saved to: {log_file}")
        log.info(f"   • Question: {len(question)} chars")
        log.info(f"   • Answer: {len(answer)} chars")
        log.info(f"   • Chunks: {len(chunks_data)}")
        log.info(f"   • Total time: {retrieval_time + generation_time:.2f}s")

        return str(log_file)

    except Exception as e:
        log.warning(f"Failed to save query log: {type(e).__name__}: {e}")
        return None


# -----------------------
# Retriever with verbose logs (this is where you "see retrieval")
# -----------------------

class HybridRetriever(BaseRetriever):
    """
    Hybrid retriever that combines:
    1. BM25 (keyword-based sparse retrieval)
    2. Vector similarity (semantic dense retrieval)

    Controlled by HYBRID_ALPHA:
    - 0.0 = pure BM25 (keyword only)
    - 0.5 = balanced hybrid
    - 1.0 = pure vector (semantic only, default)
    """

    def __init__(
        self,
        vector_store: PGVectorStore,
        embed_model: Any,
        similarity_top_k: int,
        alpha: float = 1.0,
        enable_metadata_filter: bool = False,
        mmr_threshold: float = 0.0
    ):
        self._vector_store = vector_store
        self._embed_model = embed_model
        self._similarity_top_k = similarity_top_k
        self._alpha = alpha  # Weight for vector vs BM25 (0=BM25, 1=vector)
        self._enable_metadata_filter = enable_metadata_filter
        self._mmr_threshold = mmr_threshold  # 0=disabled, >0=enable MMR diversity
        self.last_retrieval_time = 0.0

        # BM25 components (initialized on first retrieval)
        self._bm25 = None
        self._all_nodes = None
        self._tokenized_corpus = None

        super().__init__()

    def _init_bm25(self):
        """Initialize BM25 index by loading all documents from vector store."""
        if self._bm25 is not None:
            return  # Already initialized

        log.info("  🔧 Initializing BM25 index for hybrid search...")
        t = now_ms()

        try:
            # Try to import rank_bm25
            from rank_bm25 import BM25Okapi
        except ImportError:
            log.warning("  ⚠️  rank-bm25 not installed. Install with: pip install rank-bm25")
            log.warning("  ⚠️  Falling back to pure vector search")
            self._alpha = 1.0  # Force pure vector mode
            return

        # Query all nodes from vector store
        # Note: This loads all documents into memory for BM25
        vsq = VectorStoreQuery(
            query_embedding=None,
            similarity_top_k=10000,  # Get many documents for BM25 corpus
            mode="default"
        )

        # Create a dummy embedding to query all docs
        # Get dimension in a way that works for both HuggingFace and MLX
        try:
            # Try HuggingFace way
            embed_dim = self._embed_model._model.get_sentence_embedding_dimension()
        except (AttributeError, TypeError):
            # Fallback: use a small test embedding to get dimension
            test_emb = self._embed_model.get_text_embedding("test")
            embed_dim = len(test_emb)

        dummy_emb = [0.0] * embed_dim
        vsq.query_embedding = dummy_emb

        try:
            res = self._vector_store.query(vsq)
            self._all_nodes = res.nodes
            log.info(f"  ✓ Loaded {len(self._all_nodes)} documents for BM25 corpus")
        except Exception as e:
            log.warning(f"  ⚠️  Failed to load corpus for BM25: {e}")
            log.warning("  ⚠️  Falling back to pure vector search")
            self._alpha = 1.0
            return

        # Tokenize corpus for BM25
        self._tokenized_corpus = [doc.get_content().lower().split() for doc in self._all_nodes]
        self._bm25 = BM25Okapi(self._tokenized_corpus)

        log.info(f"  ✓ BM25 index initialized in {dur_s(t):.2f}s")

    def _retrieve(self, query_bundle: QueryBundle) -> List[NodeWithScore]:
        q = query_bundle.query_str
        retrieval_start = now_ms()

        log.info(f"\n{'='*70}")
        if self._alpha == 1.0:
            log.info(f"STEP 4: RETRIEVAL - Pure vector similarity search")
        elif self._alpha == 0.0:
            log.info(f"STEP 4: RETRIEVAL - Pure BM25 keyword search")
        else:
            log.info(f"STEP 4: HYBRID RETRIEVAL - Combining BM25 + Vector (α={self._alpha})")
        log.info(f"{'='*70}")
        log.info(f"❓ Query: \"{q}\"")

        # Parse metadata filters from query if enabled
        metadata_filters = {}
        filtered_query = q

        if self._enable_metadata_filter:
            metadata_filters, filtered_query = self._parse_metadata_filters(q)
            if metadata_filters:
                log.info(f"\n🔍 Metadata filters detected:")
                for key, value in metadata_filters.items():
                    log.info(f"  • {key}: {value}")
                log.info(f"  • Filtered query: \"{filtered_query}\"")

        # Vector similarity retrieval
        vector_results = []
        if self._alpha > 0:
            log.info(f"\n💡 Vector similarity search (weight={self._alpha:.1f}):")
            t0 = now_ms()
            q_emb = self._embed_model.get_query_embedding(filtered_query)
            log.info(f"  • Query embedding computed in {dur_s(t0):.2f}s")

            vsq = VectorStoreQuery(
                query_embedding=q_emb,
                similarity_top_k=self._similarity_top_k * 2,  # Get more candidates for filtering
                mode="default",
            )
            res = self._vector_store.query(vsq)
            vector_results = [(node, score) for node, score in zip(res.nodes, res.similarities or [])]
            log.info(f"  • Retrieved {len(vector_results)} candidates")

        # BM25 keyword retrieval
        bm25_results = []
        if self._alpha < 1.0:
            log.info(f"\n💡 BM25 keyword search (weight={1-self._alpha:.1f}):")
            self._init_bm25()

            if self._bm25 is not None:
                t1 = now_ms()
                tokenized_query = filtered_query.lower().split()
                bm25_scores = self._bm25.get_scores(tokenized_query)

                # Get top-k by BM25 score
                import numpy as np
                top_indices = np.argsort(bm25_scores)[::-1][:self._similarity_top_k * 2]

                for idx in top_indices:
                    if idx < len(self._all_nodes):
                        bm25_results.append((self._all_nodes[idx], bm25_scores[idx]))

                log.info(f"  • Retrieved {len(bm25_results)} candidates in {dur_s(t1):.2f}s")

        # Combine results with hybrid scoring
        log.info(f"\n🔀 Combining results:")
        combined = self._combine_results(vector_results, bm25_results, self._alpha)

        # Apply metadata filtering if enabled
        if metadata_filters:
            before_filter_count = len(combined)
            combined = self._apply_metadata_filters(combined, metadata_filters)
            log.info(f"  • Before filtering: {before_filter_count} candidates")
            log.info(f"  • After filtering: {len(combined)} results")

            if len(combined) == 0 and before_filter_count > 0:
                log.warning(f"  ⚠️  All results filtered out! Check filter criteria:")
                log.warning(f"     Filters: {metadata_filters}")
                # Show sample metadata from first few candidates for debugging
                sample_nodes = [node for node, _ in combined[:3]] if combined else [node for node, _ in [(n, s) for n, s in zip(vector_results[:3] if vector_results else [], [0]*3)]]
                if not sample_nodes and vector_results:
                    sample_nodes = [node for node, _ in vector_results[:3]]
                for node, _ in (vector_results[:3] if vector_results else []):
                    md = node.metadata or {}
                    if md.get('is_chat_log'):
                        log.warning(f"     Sample metadata: participants={md.get('participants', [])}, dates={md.get('earliest_date', '?')}-{md.get('latest_date', '?')}")
                        break

        # Apply MMR for diversity if enabled
        if self._mmr_threshold > 0 and len(combined) > self._similarity_top_k:
            log.info(f"\n🎲 Applying MMR (Maximal Marginal Relevance) for diversity:")
            log.info(f"  • Lambda: {self._mmr_threshold:.2f} (0=max diversity, 1=max relevance)")
            combined = self._apply_mmr(combined, self._similarity_top_k, self._mmr_threshold)
            log.info(f"  • Selected {len(combined)} diverse results")
        else:
            # Take top-k
            combined = combined[:self._similarity_top_k]

        # Calculate quality metrics
        scores = [score for _, score in combined if score is not None]
        if scores:
            log.info(f"\n📊 Retrieval Quality Metrics:")
            log.info(f"  • Best match score: {max(scores):.4f}")
            log.info(f"  • Worst match score: {min(scores):.4f}")
            log.info(f"  • Average score: {sum(scores)/len(scores):.4f}")

        # Log retrieved chunks
        log.info(f"\n📄 Retrieved Chunks (top {len(combined)}):")
        out = []
        for i, (node, score) in enumerate(combined, start=1):
            md = node.metadata or {}
            score_str = f"{score:.4f}" if score is not None else "None"
            page = md.get("page_label") or md.get("page") or md.get("source") or "?"

            # Show chat metadata if available
            chat_info = ""
            if md.get("is_chat_log"):
                participants = md.get("participants", [])
                date_range = md.get("date_range", "")
                if participants:
                    chat_info = f" | Participants: {', '.join(participants[:2])}"
                if date_range:
                    chat_info += f" | {date_range}"

            log.info(f"\n  {i}. Score: {score_str} | Source: {page}{chat_info}")
            chunk_text = preview(node.get_content(), 200)
            colored_text = colorize_participants(chunk_text)
            log.info(f"     Text: \"{colored_text}\"")

            nws = NodeWithScore(node=node, score=score)
            out.append(nws)

        self.last_retrieval_time = dur_s(retrieval_start)
        return out

    def _combine_results(
        self,
        vector_results: List[Tuple[Any, float]],
        bm25_results: List[Tuple[Any, float]],
        alpha: float
    ) -> List[Tuple[Any, float]]:
        """Combine vector and BM25 results with weighted scoring."""
        from collections import defaultdict

        # Normalize scores to [0, 1]
        def normalize_scores(results):
            if not results:
                return {}
            scores = [score for _, score in results]
            min_score = min(scores)
            max_score = max(scores)
            range_score = max_score - min_score if max_score > min_score else 1.0

            normalized = {}
            for node, score in results:
                node_id = node.node_id if hasattr(node, 'node_id') else id(node)
                normalized[node_id] = (node, (score - min_score) / range_score)
            return normalized

        vector_norm = normalize_scores(vector_results)
        bm25_norm = normalize_scores(bm25_results)

        # Combine scores
        all_node_ids = set(vector_norm.keys()) | set(bm25_norm.keys())
        combined = []

        for node_id in all_node_ids:
            node = None
            vector_score = 0.0
            bm25_score = 0.0

            if node_id in vector_norm:
                node, vector_score = vector_norm[node_id]
            if node_id in bm25_norm:
                node, bm25_score = bm25_norm[node_id]

            # Weighted combination
            hybrid_score = alpha * vector_score + (1 - alpha) * bm25_score
            combined.append((node, hybrid_score))

        # Sort by hybrid score
        combined.sort(key=lambda x: x[1], reverse=True)
        return combined

    def _parse_metadata_filters(self, query: str) -> Tuple[dict, str]:
        """
        Parse metadata filters from query.
        Supports: participant:Name, date:YYYY-MM-DD, after:YYYY-MM-DD, before:YYYY-MM-DD
        """
        import re

        filters = {}
        filtered_query = query

        # Extract participant filter: participant:Name or from:Name
        participant_pattern = r'(?:participant|from):([^\s]+)'
        participant_match = re.search(participant_pattern, query, re.IGNORECASE)
        if participant_match:
            filters['participant'] = participant_match.group(1)
            filtered_query = re.sub(participant_pattern, '', filtered_query, flags=re.IGNORECASE)

        # Extract date filters
        date_pattern = r'(?:after|since):(\d{4}-\d{2}-\d{2})'
        date_match = re.search(date_pattern, query, re.IGNORECASE)
        if date_match:
            filters['after'] = date_match.group(1)
            filtered_query = re.sub(date_pattern, '', filtered_query, flags=re.IGNORECASE)

        before_pattern = r'before:(\d{4}-\d{2}-\d{2})'
        before_match = re.search(before_pattern, query, re.IGNORECASE)
        if before_match:
            filters['before'] = before_match.group(1)
            filtered_query = re.sub(before_pattern, '', filtered_query, flags=re.IGNORECASE)

        # Clean up extra whitespace
        filtered_query = ' '.join(filtered_query.split())

        return filters, filtered_query

    def _apply_metadata_filters(
        self,
        results: List[Tuple[Any, float]],
        filters: dict
    ) -> List[Tuple[Any, float]]:
        """Filter results based on metadata."""
        filtered = []

        for node, score in results:
            md = node.metadata or {}

            # Check participant filter (fuzzy matching)
            if 'participant' in filters:
                participants = md.get('participants', [])
                filter_name = filters['participant'].lower()

                # Check if filter matches any participant (case-insensitive, substring match)
                matched = False
                for p in participants:
                    if filter_name in p.lower() or p.lower() in filter_name:
                        matched = True
                        break

                if not matched:
                    continue

            # Check date filters
            if 'after' in filters:
                latest_date = md.get('latest_date', '')
                if not latest_date or latest_date < filters['after']:
                    continue

            if 'before' in filters:
                earliest_date = md.get('earliest_date', '')
                if not earliest_date or earliest_date > filters['before']:
                    continue

            filtered.append((node, score))

        return filtered

    def _apply_mmr(
        self,
        results: List[Tuple[Any, float]],
        k: int,
        lambda_param: float
    ) -> List[Tuple[Any, float]]:
        """
        Apply Maximal Marginal Relevance for diversity.

        MMR balances relevance and diversity:
        - Select documents that are relevant to query
        - But also diverse from already-selected documents

        Args:
            results: List of (node, score) tuples sorted by relevance
            k: Number of results to return
            lambda_param: Trade-off between relevance and diversity
                         0.0 = max diversity, 1.0 = max relevance

        Returns:
            List of k diverse (node, score) tuples
        """
        if len(results) <= k:
            return results[:k]

        # Extract embeddings for all candidate nodes
        try:
            import numpy as np

            nodes = [node for node, _ in results]
            scores = [score for _, score in results]

            # Get embeddings (they're already stored in nodes from indexing)
            embeddings = []
            for node in nodes:
                if hasattr(node, 'embedding') and node.embedding is not None:
                    embeddings.append(node.embedding)
                else:
                    # Fallback: compute embedding on-the-fly
                    text = node.get_content()
                    emb = self._embed_model.get_text_embedding(text)
                    embeddings.append(emb)

            embeddings = np.array(embeddings)

            # Normalize embeddings for cosine similarity
            norms = np.linalg.norm(embeddings, axis=1, keepdims=True)
            embeddings = embeddings / (norms + 1e-10)

            # MMR algorithm
            selected_indices = []
            remaining_indices = list(range(len(nodes)))

            # Start with the most relevant document
            selected_indices.append(0)
            remaining_indices.remove(0)

            # Iteratively select diverse documents
            for _ in range(k - 1):
                if not remaining_indices:
                    break

                mmr_scores = []
                selected_embs = embeddings[selected_indices]

                for idx in remaining_indices:
                    candidate_emb = embeddings[idx]
                    relevance_score = scores[idx]

                    # Calculate max similarity to already-selected documents
                    similarities = np.dot(selected_embs, candidate_emb)
                    max_similarity = np.max(similarities) if len(similarities) > 0 else 0.0

                    # MMR score: balance relevance and diversity
                    mmr_score = lambda_param * relevance_score - (1 - lambda_param) * max_similarity
                    mmr_scores.append((idx, mmr_score))

                # Select document with highest MMR score
                best_idx, _ = max(mmr_scores, key=lambda x: x[1])
                selected_indices.append(best_idx)
                remaining_indices.remove(best_idx)

            # Return selected documents with their scores
            return [(nodes[idx], scores[idx]) for idx in selected_indices]

        except Exception as e:
            log.warning(f"  ⚠️  MMR failed: {e}. Falling back to top-k selection")
            return results[:k]


class VectorDBRetriever(BaseRetriever):
    """
    A retriever is "query -> relevant nodes".
    We implement it manually so you see the real steps:
      - embed query text
      - vector similarity search in Postgres
      - return NodeWithScore objects
    """

    def __init__(self, vector_store: PGVectorStore, embed_model: Any, similarity_top_k: int):
        self._vector_store = vector_store
        self._embed_model = embed_model
        self._similarity_top_k = similarity_top_k
        self.last_retrieval_time = 0.0  # Track last retrieval time for logging
        super().__init__()

    def _retrieve(self, query_bundle: QueryBundle) -> List[NodeWithScore]:
        q = query_bundle.query_str
        retrieval_start = now_ms()  # Track total retrieval time

        log.info(f"\n{'='*70}")
        log.info(f"STEP 4: RETRIEVAL - Finding relevant chunks via vector similarity")
        log.info(f"{'='*70}")
        log.info(f"❓ Query: \"{q}\"")

        log.info(f"\n💡 How retrieval works:")
        log.info(f"  1. Convert query to embedding vector (same model as documents)")
        log.info(f"  2. Calculate cosine similarity with all stored vectors")
        log.info(f"  3. Return top-{self._similarity_top_k} most similar chunks")
        log.info(f"  4. Similarity score: 1.0 = identical, 0.0 = unrelated")

        t0 = now_ms()
        q_emb = self._embed_model.get_query_embedding(q)
        log.info(f"\n🔢 Query embedding computed in {dur_s(t0):.2f}s ({len(q_emb)} dimensions)")

        t1 = now_ms()
        vsq = VectorStoreQuery(
            query_embedding=q_emb,
            similarity_top_k=self._similarity_top_k,
            mode="default",
        )
        res = self._vector_store.query(vsq)
        search_time = dur_s(t1)
        log.info(f"🔍 Vector search complete in {search_time:.2f}s")
        log.info(f"  • Searched through stored embeddings")
        log.info(f"  • Found top-{self._similarity_top_k} most similar chunks")

        out: List[NodeWithScore] = []
        for i, node in enumerate(res.nodes):
            score: Optional[float] = res.similarities[i] if res.similarities is not None else None
            nws = NodeWithScore(node=node, score=score)
            out.append(nws)

        # Calculate score distribution for quality insights
        scores = [nws.score for nws in out if isinstance(nws.score, (int, float))]
        if scores:
            avg_score = sum(scores) / len(scores)
            max_score = max(scores)
            min_score = min(scores)
            score_range = max_score - min_score

            log.info(f"\n📊 Retrieval Quality Metrics:")
            log.info(f"  • Best match score: {max_score:.4f}")
            log.info(f"  • Worst match score: {min_score:.4f}")
            log.info(f"  • Average score: {avg_score:.4f}")
            log.info(f"  • Score range: {score_range:.4f}")

            if max_score > 0.8:
                log.info(f"  ✓ Excellent: Found highly relevant chunks (>0.8)")
            elif max_score > 0.6:
                log.info(f"  ✓ Good: Found relevant chunks (0.6-0.8)")
            elif max_score > 0.4:
                log.info(f"  ⚠️  Fair: Moderate relevance (0.4-0.6) - answer may be vague")
            else:
                log.info(f"  ⚠️  Poor: Low relevance (<0.4) - may not answer question")

        # Log retrieved chunks with details
        log.info(f"\n📄 Retrieved Chunks (these will be sent to LLM):")

        # Check if retrieved chunks have configuration metadata
        chunk_configs_seen = set()
        for i, nws in enumerate(out, start=1):
            md = nws.node.metadata or {}
            score_str = f"{nws.score:.4f}" if isinstance(nws.score, (int, float)) else "None"
            page = md.get("page_label") or md.get("page") or md.get("source") or "?"

            # Track chunk configurations
            chunk_config = md.get("_index_signature", "unknown")
            chunk_configs_seen.add(chunk_config)

            # Show configuration for first chunk or if it differs
            config_info = ""
            if i == 1 or len(chunk_configs_seen) > 1:
                chunk_size = md.get("_chunk_size", "?")
                chunk_overlap = md.get("_chunk_overlap", "?")
                config_info = f" [cs={chunk_size}, ov={chunk_overlap}]"

            log.info(f"\n  {i}. Similarity: {score_str} | Source: {page}{config_info}")
            chunk_text = preview(nws.node.get_content(), 200)
            colored_text = colorize_participants(chunk_text)
            log.info(f"     Text: \"{colored_text}\"")

        # Warn if mixed configurations detected in retrieval
        if len(chunk_configs_seen) > 1:
            log.warning(f"\n  ⚠️  WARNING: Retrieved chunks from {len(chunk_configs_seen)} different configurations!")
            log.warning(f"  ⚠️  This indicates a mixed index - results may be unpredictable")
            log.warning(f"  ⚠️  Consider rebuilding with RESET_TABLE=1")

        # Store total retrieval time for logging
        self.last_retrieval_time = dur_s(retrieval_start)

        return out


# -----------------------
# Main pipeline
# -----------------------
def build_embed_model():
    """
    Build embedding model with backend selection.

    EMBED_BACKEND options:
    - huggingface (default): PyTorch + MPS/CUDA/CPU
    - mlx: Apple Silicon optimized (5-20x faster on M1/M2/M3)

    Embeddings transform text -> vector for semantic search.
    """
    backend = S.embed_backend.lower()

    # Try MLX backend first if requested
    if backend == "mlx":
        try:
            from utils.mlx_embedding import MLXEmbedding
            log.info(f"Embedding backend: MLX (Apple Silicon Metal GPU)")
            log.info(f"Embedding model: {S.embed_model_name} (dim={S.embed_dim})")
            t = now_ms()
            model = MLXEmbedding(model_name=S.embed_model_name)
            log.info(f"MLX model loaded in {dur_s(t):.2f}s")
            log.info(f"  ⚡ Expected speedup: 5-20x vs PyTorch")
            log.info(f"  💡 Tip: Increase EMBED_BATCH to 64-128 for best performance")
            return model
        except ImportError as e:
            log.warning(f"MLX not available: {e}")
            log.warning("  Install with: pip install mlx mlx-embedding-models")
            log.warning("  Falling back to HuggingFace backend")
            # Fall through to HuggingFace
        except Exception as e:
            log.error(f"MLX initialization failed: {e}")
            log.warning("  Falling back to HuggingFace backend")
            # Fall through to HuggingFace

    # HuggingFace backend (default)
    import torch
    if torch.backends.mps.is_available():
        device = "mps"  # Apple Metal GPU
        log.info(f"Embedding backend: HuggingFace with MPS (Apple Metal GPU)")
    elif torch.cuda.is_available():
        device = "cuda"
        log.info(f"Embedding backend: HuggingFace with CUDA GPU")
    else:
        device = "cpu"
        log.info(f"Embedding backend: HuggingFace with CPU")

    log.info(f"Embedding model: {S.embed_model_name} (expected dim={S.embed_dim})")
    t = now_ms()
    model = HuggingFaceEmbedding(model_name=S.embed_model_name, device=device)
    log.info(f"Embedding model loaded in {dur_s(t):.2f}s")
    return model


def build_llm():
    """
    LLM is only for "answer synthesis".
    Retrieval quality is dominated by embedding/chunking/top_k;
    LLM mainly affects answer style + reasoning and speed.

    Supports three backends (in order of preference):
    1. vLLM Server (BEST): OpenAI-compatible API, no reload (set USE_VLLM=1)
    2. vLLM Direct: In-process, reloads each time (set USE_VLLM=1, no server)
    3. llama.cpp: GGUF models, CPU/GPU (default)

    vLLM is 10-20x faster on GPU but requires HuggingFace format models.
    """
    use_vllm = os.getenv("USE_VLLM", "0") == "1"

    # Priority 1: Try vLLM Server Mode (fastest, no reload)
    if use_vllm and VLLM_CLIENT_AVAILABLE:
        try:
            vllm_model = os.getenv("VLLM_MODEL", "TheBloke/Mistral-7B-Instruct-v0.2-AWQ")
            log.info(f"LLM: vLLM Server Mode ({vllm_model})")
            log.info(f"LLM params: MAX_TOKENS={S.max_new_tokens} TEMP={S.temperature}")
            log.info("  🚀 Using vLLM server - ultra-fast queries (no model reload)!")

            llm = build_vllm_client(
                model=vllm_model,
                temperature=S.temperature,
                max_tokens=S.max_new_tokens,
            )
            return llm
        except Exception as e:
            log.warning(f"vLLM server not available: {e}")
            log.warning("Falling back to vLLM direct mode (will reload model)...")

    # Priority 2: vLLM Direct Mode (reloads each time, but still GPU-fast)
    if use_vllm and VLLM_AVAILABLE:
        vllm_model = os.getenv("VLLM_MODEL", "TheBloke/Mistral-7B-Instruct-v0.2-AWQ")
        log.info(f"LLM: vLLM Direct GPU-accelerated ({vllm_model})")
        log.info(f"LLM params: MAX_LEN={S.context_window} MAX_TOKENS={S.max_new_tokens} TEMP={S.temperature}")
        log.info("  ⚡ Using vLLM direct mode (10-20x faster than CPU)")
        log.info("  💡 For even faster queries, run vLLM server: ./scripts/start_vllm_server.sh")

        llm = build_vllm_llm(
            model_name=vllm_model,
            temperature=S.temperature,
            max_tokens=S.max_new_tokens,
            gpu_memory_utilization=0.8,
            max_model_len=S.context_window,
        )
        return llm
    elif use_vllm and not VLLM_AVAILABLE:
        log.warning("USE_VLLM=1 but vLLM not installed. Install with: pip install vllm")
        log.warning("Falling back to llama.cpp...")

    # llama.cpp backend (default, supports GGUF models)
    src = f"MODEL_PATH={S.model_path}" if S.model_path else f"MODEL_URL={S.model_url}"
    log.info(f"LLM: llama.cpp GGUF ({src})")
    log.info(f"LLM params: CTX={S.context_window} MAX_NEW_TOKENS={S.max_new_tokens} TEMP={S.temperature} N_GPU_LAYERS={S.n_gpu_layers} N_BATCH={S.n_batch}")

    llm = LlamaCPP(
        model_url=None if S.model_path else S.model_url,
        model_path=S.model_path or None,
        temperature=S.temperature,
        max_new_tokens=S.max_new_tokens,
        context_window=S.context_window,
        model_kwargs={
            "n_gpu_layers": S.n_gpu_layers,
            "n_batch": S.n_batch,
        },
        verbose=True,  # llama.cpp will emit its own logs too
    )
    return llm


def load_documents(doc_path: str) -> List[Any]:
    """
    Load documents from various formats (PDF, DOCX, TXT, MD) or from a folder.
    Returns LlamaIndex Documents.
    """
    path = Path(doc_path)

    if not path.exists():
        raise FileNotFoundError(
            f"Document not found: {doc_path}\n"
            f"  Fix: Place your document at this path or set PDF_PATH environment variable\n"
            f"  Example: PDF_PATH=/path/to/your/document.pdf python {sys.argv[0]}"
        )

    t = now_ms()
    docs = []

    # Supported text-based extensions (loaded as plain text)
    text_extensions = {
        ".txt", ".md", ".html", ".htm", ".json", ".csv", ".xml", ".xsl",
        ".py", ".js", ".ts", ".jsx", ".tsx", ".cpp", ".c", ".h", ".hpp",
        ".java", ".rb", ".go", ".rs", ".php", ".m", ".swift", ".kt",
        ".sh", ".bash", ".zsh", ".yaml", ".yml", ".toml", ".ini", ".cfg",
        ".sql", ".r", ".scala", ".pl", ".lua", ".ex", ".exs", ".clj",
    }

    # Handle directory (folder of documents)
    if path.is_dir():
        log.info(f"Loading documents from folder: {doc_path}")

        try:
            from llama_index.core import SimpleDirectoryReader
        except ImportError:
            raise ImportError(
                "SimpleDirectoryReader not available. Update llama-index: pip install -U llama-index"
            )

        # Count files first
        supported_extensions = [".pdf", ".docx", ".pptx"] + list(text_extensions)
        file_count = sum(1 for f in path.rglob("*") if f.is_file() and f.suffix.lower() in supported_extensions)

        if file_count == 0:
            raise ValueError(
                f"No supported documents found in folder: {doc_path}\n"
                f"  Supported formats: {', '.join(supported_extensions)}"
            )

        log.info(f"  Found {file_count} supported file(s)")

        # Load all documents recursively
        reader = SimpleDirectoryReader(
            input_dir=str(path),
            recursive=True,
            required_exts=supported_extensions,
        )
        docs = reader.load_data()

        # Clean HTML documents
        html_count = 0
        for doc in docs:
            source = doc.metadata.get("file_path", "") or doc.metadata.get("source", "")
            if source.lower().endswith((".html", ".htm")):
                cleaned_text = clean_html_content(doc.text)
                doc.set_content(cleaned_text)
                html_count += 1

        if html_count > 0:
            log.info(f"  🧹 Cleaned {html_count} HTML file(s) (removed tags/scripts/styles)")

        # Calculate total text statistics for folder
        total_chars = sum(len(doc.text) for doc in docs)
        total_words = sum(len(doc.text.split()) for doc in docs)
        avg_chars_per_doc = total_chars / len(docs) if docs else 0

        log.info(f"✓ Loaded {len(docs)} document(s) from folder in {dur_s(t):.2f}s")
        log.info(f"  📊 Statistics:")
        log.info(f"    • Total characters: {total_chars:,}")
        log.info(f"    • Total words: {total_words:,}")
        log.info(f"    • Avg chars/doc: {avg_chars_per_doc:.0f}")
        log.info(f"  ℹ️  Note: Loaded recursively from folder with {file_count} source files")

    else:
        # Handle single file
        ext = path.suffix.lower()
        file_size_mb = path.stat().st_size / (1024 * 1024)

        log.info(f"Loading document: {doc_path}")
        log.info(f"  Format: {ext}, Size: {file_size_mb:.1f} MB")

        try:
            if ext == ".pdf":
                # PDF: Use PyMuPDFReader (page-per-document)
                docs = PyMuPDFReader().load(file_path=str(path))

            elif ext == ".docx":
                # DOCX: Load as single document
                try:
                    from docx import Document as DocxDocument
                except ImportError:
                    raise ImportError(
                        "python-docx not installed. Install with: pip install python-docx"
                    )

                docx_doc = DocxDocument(str(path))
                text = "\n".join([p.text for p in docx_doc.paragraphs])

                from llama_index.core.schema import Document
                docs = [Document(text=text, metadata={"source": str(path), "format": "docx"})]

            elif ext == ".pptx":
                # PPTX: PowerPoint - extract text from slides
                try:
                    from pptx import Presentation
                except ImportError:
                    raise ImportError(
                        "python-pptx not installed. Install with: pip install python-pptx"
                    )

                prs = Presentation(str(path))
                text_parts = []
                for slide_num, slide in enumerate(prs.slides, 1):
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text)
                    if slide_text:
                        text_parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_text))

                text = "\n\n".join(text_parts)
                from llama_index.core.schema import Document
                docs = [Document(text=text, metadata={"source": str(path), "format": "pptx"})]

            elif ext in (".html", ".htm"):
                # HTML: Clean and extract text content
                raw_html = path.read_text(encoding="utf-8", errors="replace")
                text = clean_html_content(raw_html)

                from llama_index.core.schema import Document
                docs = [Document(text=text, metadata={"source": str(path), "format": "html"})]

            elif ext in text_extensions:
                # Text-based files: Load as single document
                text = path.read_text(encoding="utf-8", errors="replace")

                from llama_index.core.schema import Document
                docs = [Document(text=text, metadata={"source": str(path), "format": ext[1:]})]

            else:
                raise ValueError(
                    f"Unsupported file format: {ext}\n"
                    f"  Supported: .pdf, .docx, .pptx, .txt, .md, .html, .json, .csv, .xml,\n"
                    f"             .py, .js, .cpp, .java, .go, .rs, .php, .m, and more\n"
                    f"  File: {doc_path}"
                )

        except Exception as e:
            if isinstance(e, (FileNotFoundError, ValueError, ImportError)):
                raise
            raise RuntimeError(
                f"Failed to load document: {doc_path}\n"
                f"  Error: {type(e).__name__}: {e}\n"
                f"  Fix: Ensure the file is valid and not corrupted"
            )

        if not docs:
            raise ValueError(
                f"Document loaded but contains no content: {doc_path}\n"
                f"  The document may be empty or corrupted"
            )

        # Calculate total text statistics
        total_chars = sum(len(doc.text) for doc in docs)
        total_words = sum(len(doc.text.split()) for doc in docs)
        avg_chars_per_doc = total_chars / len(docs) if docs else 0

        log.info(f"✓ Loaded {len(docs)} document(s) in {dur_s(t):.2f}s")
        log.info(f"  📊 Statistics:")
        log.info(f"    • Total characters: {total_chars:,}")
        log.info(f"    • Total words: {total_words:,}")
        log.info(f"    • Avg chars/doc: {avg_chars_per_doc:.0f}")
        if ext == ".pdf":
            log.info(f"  ℹ️  Note: PDFs are split per-page for better citation accuracy")
        else:
            log.info(f"  ℹ️  Note: {ext.upper()} loaded as single document for context preservation")

    return docs


def chunk_documents(docs: List[Any]) -> Tuple[List[str], List[int]]:
    """
    Turn documents into text chunks.
    We keep doc_idxs so each chunk knows which doc/page it came from (for metadata/citations).
    """
    log.info(f"\n{'='*70}")
    log.info(f"STEP 1: CHUNKING - Breaking documents into retrievable pieces")
    log.info(f"{'='*70}")
    log.info(f"⚙️  Configuration:")
    log.info(f"  • chunk_size={S.chunk_size} chars")
    log.info(f"  • chunk_overlap={S.chunk_overlap} chars ({S.chunk_overlap/S.chunk_size*100:.1f}% overlap)")

    overlap_ratio = S.chunk_overlap / S.chunk_size
    if overlap_ratio < 0.1:
        log.info(f"  ⚠️  Low overlap may break context across chunks")
    elif overlap_ratio > 0.3:
        log.info(f"  ℹ️  High overlap preserves context but increases storage")
    else:
        log.info(f"  ✓ Good overlap ratio for context preservation")

    log.info(f"\n💡 Why chunking matters:")
    log.info(f"  • Smaller chunks = more precise retrieval, less context")
    log.info(f"  • Larger chunks = more context, less precise matching")
    log.info(f"  • Overlap prevents splitting important information")

    splitter = SentenceSplitter(chunk_size=S.chunk_size, chunk_overlap=S.chunk_overlap)

    t = now_ms()
    chunks: List[str] = []
    doc_idxs: List[int] = []
    chunk_sizes: List[int] = []

    for doc_idx, doc in enumerate(docs):
        # Split this page's text into chunks
        cs = splitter.split_text(doc.text)
        chunks.extend(cs)
        doc_idxs.extend([doc_idx] * len(cs))
        chunk_sizes.extend([len(c) for c in cs])

        # Lightweight progress log every ~25 pages
        if (doc_idx + 1) % 25 == 0:
            log.info(f"  📄 Processed {doc_idx+1}/{len(docs)} documents → {len(chunks)} chunks so far")

    avg_chunk_size = sum(chunk_sizes) / len(chunk_sizes) if chunk_sizes else 0
    min_chunk_size = min(chunk_sizes) if chunk_sizes else 0
    max_chunk_size = max(chunk_sizes) if chunk_sizes else 0

    log.info(f"\n✓ Chunking complete in {dur_s(t):.2f}s")
    log.info(f"  📊 Chunk Statistics:")
    log.info(f"    • Total chunks: {len(chunks)}")
    log.info(f"    • Avg chunk size: {avg_chunk_size:.0f} chars")
    log.info(f"    • Min chunk size: {min_chunk_size} chars")
    log.info(f"    • Max chunk size: {max_chunk_size} chars")
    log.info(f"    • Chunks per document: {len(chunks)/len(docs):.1f} average")

    if chunks:
        log.info(f"\n  📝 Example chunk (first):")
        example_text = preview(chunks[0], 150)
        colored_example = colorize_participants(example_text)
        log.info(f"    \"{colored_example}\"")

    return chunks, doc_idxs


def build_nodes(docs: List[Any], chunks: List[str], doc_idxs: List[int]) -> List[TextNode]:
    """
    Create TextNode objects (text + metadata).
    Embedding is added later.

    IMPORTANT: We store chunking parameters in metadata so we can:
    - Detect mixed-index scenarios (different chunk_size in same table)
    - Debug retrieval quality issues
    - Track what configuration produced each chunk
    """
    log.info("Building TextNode objects (text + metadata)")
    t = now_ms()

    # Create an index signature to track chunking configuration
    index_signature = f"cs{S.chunk_size}_ov{S.chunk_overlap}_{S.embed_model_name.replace('/', '_')}"

    nodes: List[TextNode] = []
    chat_logs_detected = 0

    for i, chunk in enumerate(chunks):
        n = TextNode(text=chunk)

        # Start with source document metadata
        src_doc = docs[doc_idxs[i]]
        n.metadata = src_doc.metadata.copy() if src_doc.metadata else {}

        # Add chunking parameters to metadata
        # This allows detection of mixed indexes and better debugging
        n.metadata["_chunk_size"] = S.chunk_size
        n.metadata["_chunk_overlap"] = S.chunk_overlap
        n.metadata["_embed_model"] = S.embed_model_name
        n.metadata["_index_signature"] = index_signature

        # Extract chat metadata if this looks like a chat log
        if os.getenv("EXTRACT_CHAT_METADATA", "1") == "1":
            chat_meta = extract_chat_metadata(chunk)
            if chat_meta:
                n.metadata.update(chat_meta)
                chat_logs_detected += 1

        nodes.append(n)

        if (i + 1) % 500 == 0:
            log.info(f"  built {i+1}/{len(chunks)} nodes")

    log.info(f"Built {len(nodes)} nodes in {dur_s(t):.2f}s")
    log.info(f"Index signature: {index_signature}")

    if chat_logs_detected > 0:
        log.info(f"✓ Detected {chat_logs_detected} chat log chunks with metadata extracted")
        log.info(f"  • Metadata includes: participants, dates, message counts")
        log.info(f"  • This enables advanced filtering and hybrid search")

    return nodes


def embed_nodes(embed_model: HuggingFaceEmbedding, nodes: List[TextNode]) -> None:
    """
    Compute embeddings for each node.
    This is often the longest step after LLM inference.
    We do batching for speed and steadier memory usage.
    """
    log.info(f"\n{'='*70}")
    log.info(f"STEP 2: EMBEDDING - Converting text to semantic vectors")
    log.info(f"{'='*70}")
    log.info(f"⚙️  Configuration:")
    log.info(f"  • Model: {S.embed_model_name}")
    log.info(f"  • Dimension: {S.embed_dim} (vector size)")
    log.info(f"  • Batch size: {S.embed_batch} chunks/batch")
    log.info(f"  • Total nodes: {len(nodes)}")

    log.info(f"\n💡 What are embeddings?")
    log.info(f"  • Embeddings convert text into {S.embed_dim}-dimensional vectors")
    log.info(f"  • Similar text → similar vectors (measured by cosine similarity)")
    log.info(f"  • This enables semantic search (meaning-based, not just keywords)")
    log.info(f"  • Example: 'cat' and 'feline' have similar embeddings")

    total_bytes = len(nodes) * S.embed_dim * 4  # 4 bytes per float32
    log.info(f"\n📊 Storage Impact:")
    log.info(f"  • {len(nodes)} vectors × {S.embed_dim} dims × 4 bytes = {total_bytes/1024/1024:.1f} MB")

    t = now_ms()
    total = len(nodes)

    # We embed only the node text (metadata_mode="none") to keep embeddings "pure".
    # You can switch to metadata_mode="all" to include metadata in embeddings (sometimes helps, sometimes adds noise).
    texts = [n.get_content(metadata_mode="none") for n in nodes]

    done = 0
    batches = list(chunked(texts, S.embed_batch))

    log.info(f"\n🔄 Processing {len(batches)} batches...")
    log.info("")  # Extra line for better progress bar visibility

    # Use tqdm if available for better progress visualization
    if tqdm:
        iterator = tqdm(
            enumerate(batches, start=1),
            total=len(batches),
            desc="⚡ Embedding batches",
            unit="batch",
            ncols=100,
            position=0,
            leave=True,
            bar_format="{l_bar}{bar}| {n_fmt}/{total_fmt} [{elapsed}<{remaining}, {rate_fmt}]"
        )
    else:
        iterator = enumerate(batches, start=1)

    batch_times = []
    for batch_idx, batch_texts in iterator:
        tb = now_ms()

        # Newer versions of HuggingFaceEmbedding support batch embedding
        # If not, this will throw and we fallback to per-item embedding.
        try:
            batch_embs = embed_model.get_text_embedding_batch(batch_texts)
        except Exception:
            batch_embs = [embed_model.get_text_embedding(x) for x in batch_texts]

        # Write embeddings back into the nodes (must align exactly)
        start_i = (batch_idx - 1) * S.embed_batch
        for j, emb in enumerate(batch_embs):
            nodes[start_i + j].embedding = emb

        done += len(batch_texts)
        batch_time = dur_s(tb)
        batch_times.append(batch_time)

        # Update progress bar with throughput info
        if tqdm:
            rate = done / max(dur_s(t), 1e-6)
            iterator.set_postfix({"nodes/s": f"{rate:.1f}", "nodes": f"{done}/{total}"})
        else:
            # Progress logs without tqdm
            rate = done / max(dur_s(t), 1e-6)
            log.info(f"  📦 Batch {batch_idx:04d} → {done}/{total} nodes | {batch_time:.2f}s | ~{rate:.1f} nodes/s")

    total_time = dur_s(t)
    avg_batch_time = sum(batch_times) / len(batch_times) if batch_times else 0
    throughput = total / total_time if total_time > 0 else 0

    log.info(f"\n✓ Embeddings complete in {total_time:.2f}s")
    log.info(f"  📊 Performance:")
    log.info(f"    • Average batch time: {avg_batch_time:.2f}s")
    log.info(f"    • Throughput: {throughput:.1f} nodes/second")
    log.info(f"    • Time per node: {total_time/total*1000:.1f}ms")
    log.info(f"\n  ℹ️  These embeddings enable semantic similarity search in the next step")


def make_vector_store() -> PGVectorStore:
    """
    Create the vector store client that uses Postgres + pgvector.
    """
    log.info(f"Connecting to Postgres vector store: db={S.db_name} host={S.host}:{S.port} user={S.user} table={S.table}")
    store = PGVectorStore.from_params(
        database=S.db_name,
        host=S.host,
        port=S.port,
        user=S.user,
        password=S.password,
        table_name=S.table,
        embed_dim=S.embed_dim,
    )
    return store


def insert_nodes(vector_store: PGVectorStore, nodes: List[TextNode]) -> None:
    """
    Insert nodes into Postgres.
    Depending on LlamaIndex version, PGVectorStore.add may do its own batching internally.
    We'll still batch to:
      - keep transactions smaller
      - show progress clearly
    """
    log.info(f"\n{'='*70}")
    log.info(f"STEP 3: STORAGE - Saving embeddings to PostgreSQL + pgvector")
    log.info(f"{'='*70}")

    before = count_rows()
    if before is not None:
        log.info(f"📊 Current table state: {before} rows exist")

    batch_size = int(os.getenv("DB_INSERT_BATCH", "250"))
    log.info(f"⚙️  Configuration:")
    log.info(f"  • Table: {S.table}")
    log.info(f"  • Batch size: {batch_size} nodes/batch")
    log.info(f"  • Total to insert: {len(nodes)} nodes")

    log.info(f"\n💡 Why use pgvector?")
    log.info(f"  • Specialized PostgreSQL extension for vector similarity search")
    log.info(f"  • Efficient indexing (IVFFlat, HNSW) for fast retrieval")
    log.info(f"  • Supports cosine distance, L2 distance, inner product")
    log.info(f"  • Production-ready: ACID compliance + SQL features")

    t = now_ms()
    total = len(nodes)
    inserted = 0

    batches = list(chunked(nodes, batch_size))
    log.info(f"\n🔄 Inserting {len(batches)} batches...")
    iterator = tqdm(enumerate(batches, start=1), total=len(batches), desc="Inserting", unit="batch") if tqdm else enumerate(batches, start=1)

    for bidx, batch in iterator:
        tb = now_ms()
        vector_store.add(batch)
        inserted += len(batch)

        if not tqdm:
            log.info(f"  💾 Batch {bidx:04d} → {inserted}/{total} nodes | {dur_s(tb):.2f}s")

    total_time = dur_s(t)
    throughput = total / total_time if total_time > 0 else 0

    after = count_rows()
    log.info(f"\n✓ Storage complete in {total_time:.2f}s")
    log.info(f"  📊 Results:")
    log.info(f"    • Inserted: {total} nodes")
    log.info(f"    • Throughput: {throughput:.1f} nodes/second")
    if after is not None:
        log.info(f"    • Total rows in table: {after}")
        if before is not None:
            log.info(f"    • New rows added: {after - before}")
    log.info(f"\n  ✓ Vector index is now ready for semantic search!")


def create_hnsw_index(table_name: str = None) -> bool:
    """
    Create an HNSW index on the embedding column for fast similarity search.

    HNSW (Hierarchical Navigable Small World) provides:
      - O(log n) search complexity vs O(n) for sequential scan
      - 50-100x faster retrieval on tables with >10K rows
      - Approximate nearest neighbor with high recall (typically >95%)

    Args:
        table_name: Table name (without 'data_' prefix). Uses S.table if not provided.

    Returns:
        True if index was created or already exists, False on error.
    """
    table = table_name or S.table
    actual_table = f"data_{table}"
    index_name = f"{actual_table}_hnsw_idx"

    log.info(f"\n{'='*70}")
    log.info(f"STEP 4: INDEXING - Creating HNSW index for fast retrieval")
    log.info(f"{'='*70}")

    log.info(f"\n💡 Why HNSW indexing matters:")
    log.info(f"  • Without index: Sequential scan compares query to EVERY embedding")
    log.info(f"  • With HNSW: Graph-based search finds neighbors in O(log n) time")
    log.info(f"  • Speedup: 50-100x faster on tables with >10,000 rows")
    log.info(f"  • Trade-off: Slightly approximate results (typically >95% recall)")

    try:
        conn = db_conn()
        conn.autocommit = True

        with conn.cursor() as cur:
            # Check if index already exists
            cur.execute("""
                SELECT indexname FROM pg_indexes
                WHERE tablename = %s AND indexname = %s
            """, (actual_table, index_name))

            if cur.fetchone():
                log.info(f"\n✓ HNSW index '{index_name}' already exists")
                conn.close()
                return True

            # Get row count to estimate indexing time
            cur.execute(f'SELECT COUNT(*) FROM "{actual_table}"')
            row_count = cur.fetchone()[0]

            log.info(f"\n⚙️  Creating HNSW index:")
            log.info(f"  • Table: {actual_table}")
            log.info(f"  • Index name: {index_name}")
            log.info(f"  • Rows to index: {row_count:,}")
            log.info(f"  • Parameters: m=16 (connections), ef_construction=64 (build quality)")

            if row_count > 50000:
                log.info(f"  ⚠️  Large table - indexing may take several minutes...")

            t = now_ms()

            # Create the HNSW index
            # m=16: Number of connections per layer (higher = better recall, more memory)
            # ef_construction=64: Build-time search width (higher = better quality, slower build)
            cur.execute(f'''
                CREATE INDEX "{index_name}"
                ON "{actual_table}"
                USING hnsw (embedding vector_cosine_ops)
                WITH (m = 16, ef_construction = 64)
            ''')

            index_time = dur_s(t)

            log.info(f"\n✓ HNSW index created in {index_time:.2f}s")
            log.info(f"  • Indexed {row_count:,} embeddings")
            log.info(f"  • Throughput: {row_count / index_time:.1f} vectors/second")
            log.info(f"\n  🚀 Similarity search is now optimized!")

        conn.close()
        return True

    except Exception as e:
        log.error(f"✗ Failed to create HNSW index: {e}")
        log.error(f"  You can create it manually with:")
        log.error(f"  CREATE INDEX ON \"{actual_table}\" USING hnsw (embedding vector_cosine_ops);")
        return False


def health_check() -> None:
    """
    Perform startup health checks to catch common issues early.
    Better to fail fast with clear errors than fail later with cryptic ones.
    """
    log.info("=== HEALTH CHECK ===")
    checks_passed = 0
    checks_total = 0

    # Check 1: PostgreSQL connectivity
    checks_total += 1
    try:
        log.info("✓ Checking PostgreSQL connection...")
        conn = admin_conn()
        conn.close()
        log.info("  ✓ PostgreSQL is reachable")
        checks_passed += 1
    except Exception as e:
        log.error(f"  ✗ PostgreSQL connection failed: {e}")
        log.error("    Fix: Ensure PostgreSQL is running (docker-compose up -d)")

    # Check 2: PDF file exists
    checks_total += 1
    try:
        log.info("✓ Checking PDF file...")
        if Path(S.pdf_path).exists():
            size_mb = os.path.getsize(S.pdf_path) / (1024 * 1024)
            log.info(f"  ✓ PDF file exists ({size_mb:.1f} MB)")
            checks_passed += 1
        else:
            log.error(f"  ✗ PDF file not found: {S.pdf_path}")
    except Exception as e:
        log.error(f"  ✗ PDF check failed: {e}")

    # Check 3: Disk space (if psutil available)
    if psutil:
        checks_total += 1
        try:
            log.info("✓ Checking disk space...")
            disk = psutil.disk_usage('/')
            free_gb = disk.free / (1024**3)
            if free_gb < 2:
                log.warning(f"  ⚠ Low disk space: {free_gb:.1f} GB free (recommend at least 5GB)")
            else:
                log.info(f"  ✓ Disk space OK ({free_gb:.1f} GB free)")
                checks_passed += 1
        except Exception as e:
            log.warning(f"  ⚠ Disk space check failed: {e}")

    # Check 4: Memory availability (if psutil available)
    if psutil:
        checks_total += 1
        try:
            log.info("✓ Checking memory...")
            vm = psutil.virtual_memory()
            avail_gb = vm.available / (1024**3)
            if avail_gb < 2:
                log.warning(f"  ⚠ Low memory: {avail_gb:.1f} GB available (recommend at least 4GB)")
                log.warning("    Consider closing other applications or reducing N_GPU_LAYERS")
            else:
                log.info(f"  ✓ Memory OK ({avail_gb:.1f} GB available)")
                checks_passed += 1
        except Exception as e:
            log.warning(f"  ⚠ Memory check failed: {e}")

    log.info(f"Health check: {checks_passed}/{checks_total} checks passed")

    if checks_passed < 2:  # At minimum need DB and PDF
        log.error("Critical health checks failed. Cannot proceed.")
        log.error("Please fix the issues above and try again.")
        sys.exit(1)
    elif checks_passed < checks_total:
        log.warning("Some health checks failed, but proceeding anyway...")
        time.sleep(2)  # Give user time to see warnings

    log.info("Health check complete!\n")


def run_query(query_engine: Any, question: str, show_sources: bool = True, retrieval_time: float = 0.0) -> None:
    """Execute a single query and display results."""
    log.info(f"\n{'='*70}")
    log.info(f"STEP 5: GENERATION - LLM synthesizes answer from retrieved chunks")
    log.info(f"{'='*70}")

    log.info(f"\n💡 How answer generation works:")
    log.info(f"  1. Retrieved chunks are combined into context")
    log.info(f"  2. Context + query sent to local LLM")
    log.info(f"  3. LLM generates answer grounded in provided context")
    log.info(f"  4. This prevents hallucination (making up information)")

    log.info(f"\n⚙️  LLM Configuration:")
    log.info(f"  • Model: Mistral 7B Instruct (Q4_K_M quantized)")
    log.info(f"  • Temperature: {S.temperature} ({'factual/deterministic' if S.temperature < 0.3 else 'balanced' if S.temperature < 0.7 else 'creative'})")
    log.info(f"  • Max tokens: {S.max_new_tokens}")
    log.info(f"  • Context window: {S.context_window}")

    log.info(f"\n🔄 Generating answer...")
    t = now_ms()
    resp = query_engine.query(question)
    generation_time = dur_s(t)

    # Calculate token statistics
    answer_length = len(str(resp))
    answer_words = len(str(resp).split())
    tokens_per_second = answer_words / generation_time if generation_time > 0 else 0

    log.info(f"\n✓ Answer generated in {generation_time:.2f}s")
    log.info(f"  📊 Generation Stats:")
    log.info(f"    • Answer length: {answer_length} characters, {answer_words} words")
    log.info(f"    • Speed: ~{tokens_per_second:.1f} words/second")
    log.info(f"    • Time per word: {generation_time/answer_words*1000:.0f}ms") if answer_words > 0 else None

    print("\n" + "="*70)
    print("✨ FINAL ANSWER:")
    print("="*70)
    print(str(resp))
    print("="*70 + "\n")

    # Show sources for transparency
    if show_sources and resp.source_nodes:
        log.info(f"\n📚 Sources Used ({len(resp.source_nodes)} chunks):")
        for i, node in enumerate(resp.source_nodes[:3], 1):  # Show top 3
            md = node.node.metadata or {}
            page = md.get("page_label") or md.get("page") or md.get("source") or "?"
            score = node.score if hasattr(node, 'score') else None
            score_str = f" (similarity: {score:.4f})" if score else ""
            log.info(f"  {i}. Source: {page}{score_str}")
            source_text = preview(node.node.get_content(), 150)
            colored_source = colorize_participants(source_text)
            log.info(f"     \"{colored_source}\"")
        if len(resp.source_nodes) > 3:
            log.info(f"  ... and {len(resp.source_nodes) - 3} more chunks")

        log.info(f"\n  ℹ️  Answer is grounded in these retrieved chunks")

    # Save query log if enabled
    if os.getenv("LOG_QUERIES", "0") == "1":
        parameters = {
            "chunk_size": S.chunk_size,
            "chunk_overlap": S.chunk_overlap,
            "top_k": S.top_k,
            "temperature": S.temperature,
            "max_new_tokens": S.max_new_tokens,
            "context_window": S.context_window,
            "n_gpu_layers": S.n_gpu_layers,
            "n_batch": S.n_batch,
            "embed_model": S.embed_model_name,
            "embed_dim": S.embed_dim,
            "embed_batch": S.embed_batch
        }

        save_query_log(
            question=question,
            answer=str(resp),
            retrieved_chunks=resp.source_nodes if hasattr(resp, 'source_nodes') else [],
            retrieval_time=retrieval_time,
            generation_time=generation_time,
            parameters=parameters
        )


def interactive_mode(query_engine: Any, retriever: VectorDBRetriever) -> None:
    """
    Interactive REPL mode for asking multiple questions.
    Type 'exit' or 'quit' to end session.
    """
    print("\n" + "="*70)
    print("INTERACTIVE MODE")
    print("="*70)
    print("Ask questions about your documents. Type 'exit' or 'quit' to end.\n")

    while True:
        try:
            question = input("Question: ").strip()

            if not question:
                continue

            if question.lower() in ('exit', 'quit', 'q'):
                print("Goodbye!")
                break

            run_query(query_engine, question, show_sources=False, retrieval_time=retriever.last_retrieval_time)

        except KeyboardInterrupt:
            print("\n\nInterrupted. Goodbye!")
            break
        except EOFError:
            print("\nGoodbye!")
            break
        except Exception as e:
            log.error(f"Query failed: {type(e).__name__}: {e}")
            print(f"Error: {e}\n")


def main():
    # Parse CLI arguments first
    args = parse_args()

    # Override settings from CLI args
    if args.doc:
        S.pdf_path = args.doc
        # Regenerate table name if PGTABLE wasn't explicitly set
        if not os.getenv("PGTABLE"):
            S.table = generate_table_name(
                S.pdf_path,
                S.chunk_size,
                S.chunk_overlap,
                S.embed_model_name
            )
            log.debug(f"Regenerated table name for CLI doc: {S.table}")
    if args.query:
        S.question = args.query

    # Validate settings (unless explicitly skipped)
    if not args.skip_validation:
        S.validate()

    log_system_info()

    log.info("=== SETTINGS ===")
    log.info(f"DB: postgresql://{S.user}:***@{S.host}:{S.port}/{S.db_name} table={S.table}")
    log.info(f"Document: {S.pdf_path}")
    log.info(f"Chunking: chunk_size={S.chunk_size} overlap={S.chunk_overlap}")
    log.info(f"Retrieval: TOP_K={S.top_k}")
    log.info(f"Embeddings: model={S.embed_model_name} dim={S.embed_dim} batch={S.embed_batch}")
    log.info(f"LLM: CTX={S.context_window} MAX_NEW_TOKENS={S.max_new_tokens} TEMP={S.temperature} N_GPU_LAYERS={S.n_gpu_layers} N_BATCH={S.n_batch}")
    mode = 'Query-only' if args.query_only else 'Index-only' if args.index_only else 'Interactive' if args.interactive else 'Full pipeline'
    log.info(f"Mode: {mode}")

    # Info about table naming
    if not os.getenv("PGTABLE"):
        log.info("")
        log.info("📋 Auto-generated table name based on configuration:")
        log.info(f"   Format: {{doc}}_cs{{size}}_ov{{overlap}}_{{model}}_{{YYMMDD}}")
        log.info(f"   To use a custom name, set: PGTABLE=my_table_name")
    else:
        log.info("")
        log.info("📋 Using explicitly set table name from PGTABLE env var")

    if not args.skip_validation:
        # Run health checks before starting heavy operations
        health_check()

    # --- DB prep ---
    ensure_db_exists()
    ensure_pgvector_extension()

    # --- Ingestion pipeline (skip if query-only mode) ---
    if not args.query_only:
        log.info("=== INDEXING DOCUMENTS ===")

        # Check if table already has data
        existing_rows = count_rows()
        if existing_rows and existing_rows > 0 and not S.reset_table:
            log.warning(f"Table '{S.table}' already contains {existing_rows} rows")

            # Check the configuration of existing index
            existing_config = check_index_configuration()
            current_signature = f"cs{S.chunk_size}_ov{S.chunk_overlap}_{S.embed_model_name.replace('/', '_')}"

            if existing_config:
                if existing_config.get("legacy") or not existing_config.get("has_metadata"):
                    log.warning("  ⚠️  LEGACY INDEX: Existing rows don't have configuration metadata")
                    log.warning("  ⚠️  Cannot detect if chunk_size matches current settings")
                elif not existing_config.get("is_consistent"):
                    log.error("  ❌ MIXED INDEX DETECTED: Table contains chunks from multiple configurations!")
                    log.error("  ❌ This will produce unreliable retrieval results")
                    log.error("  ❌ You should set RESET_TABLE=1 to clean up")
                elif existing_config.get("index_signature") != current_signature:
                    log.error("  ❌ CONFIGURATION MISMATCH DETECTED!")
                    log.error(f"  Current config: chunk_size={S.chunk_size}, overlap={S.chunk_overlap}, model={S.embed_model_name}")
                    log.error(f"  Existing index: chunk_size={existing_config.get('chunk_size')}, overlap={existing_config.get('chunk_overlap')}, model={existing_config.get('embed_model')}")
                    log.error("")
                    log.error("  ❌ Proceeding will create a MIXED INDEX (same table, different chunk sizes)")
                    log.error("  ❌ This causes retrieval to return chunks from BOTH configurations")
                    log.error("  ❌ The chunk_size parameter will appear to 'do nothing' because old chunks still exist")
                    log.error("")
                    log.error("  FIX OPTIONS:")
                    log.error("    1. Set RESET_TABLE=1 to rebuild with new config")
                    log.error(f"    2. Use different table name: PGTABLE=ethical-slut_cs{S.chunk_size}_ov{S.chunk_overlap}")
                    log.error("    3. Use --query-only to query existing index without adding more rows")
                    log.error("")
                    sys.exit(1)
                else:
                    log.info(f"  ✓ Configuration matches existing index: {existing_config.get('index_signature')}")

            log.warning("  Set RESET_TABLE=1 to re-index, or use --query-only to skip indexing")
            log.warning("  Proceeding will add MORE rows (incremental indexing)")

        reset_table_if_requested()

        # Print all RAG-related environment variables for verification
        log.info("=== ENVIRONMENT VARIABLES ===")
        rag_env_vars = [
            "CHUNK_SIZE", "CHUNK_OVERLAP", "TOP_K",
            "EMBED_MODEL", "EMBED_DIM", "EMBED_BATCH", "EMBED_BACKEND",
            "N_GPU_LAYERS", "N_BATCH", "CTX", "MAX_NEW_TOKENS", "TEMP",
            "HYBRID_ALPHA", "MMR_THRESHOLD", "ENABLE_FILTERS",
            "LOG_FULL_CHUNKS", "COLORIZE_CHUNKS", "LOG_LEVEL",
            "EXTRACT_CHAT_METADATA", "PGTABLE", "PDF_PATH", "RESET_TABLE"
        ]
        for var in rag_env_vars:
            value = os.getenv(var, "(not set)")
            log.info(f"  {var:25s} = {value}")
        log.info("=" * 70)

        # Load embedding model for indexing
        embed_model = build_embed_model()

        # --- Vector store client ---
        vector_store = make_vector_store()

        # --- Document loading and processing ---
        docs = load_documents(S.pdf_path)
        chunks, doc_idxs = chunk_documents(docs)
        nodes = build_nodes(docs, chunks, doc_idxs)

        embed_nodes(embed_model, nodes)
        insert_nodes(vector_store, nodes)
        create_hnsw_index()  # Create HNSW index for fast similarity search

        # Exit early if index-only mode (skip LLM loading and queries)
        if args.index_only:
            log.info("=== INDEXING COMPLETE ===")
            log.info(f"  Indexed {len(nodes)} chunks into table '{S.table}'")
            log.info("  Use --query-only or option 2 in interactive menu to query this index")
            return
    else:
        log.info("=== SKIPPING INDEXING (query-only mode) ===")

        # Print all RAG-related environment variables for verification
        log.info("=== ENVIRONMENT VARIABLES ===")
        rag_env_vars = [
            "CHUNK_SIZE", "CHUNK_OVERLAP", "TOP_K",
            "EMBED_MODEL", "EMBED_DIM", "EMBED_BATCH", "EMBED_BACKEND",
            "N_GPU_LAYERS", "N_BATCH", "CTX", "MAX_NEW_TOKENS", "TEMP",
            "HYBRID_ALPHA", "MMR_THRESHOLD", "ENABLE_FILTERS",
            "LOG_FULL_CHUNKS", "COLORIZE_CHUNKS", "LOG_LEVEL",
            "EXTRACT_CHAT_METADATA", "PGTABLE", "PDF_PATH", "RESET_TABLE"
        ]
        for var in rag_env_vars:
            value = os.getenv(var, "(not set)")
            log.info(f"  {var:25s} = {value}")
        log.info("=" * 70)

        # Still need embed model for queries
        embed_model = build_embed_model()
        vector_store = make_vector_store()

        # Verify table exists and has data
        existing_rows = count_rows()
        if not existing_rows or existing_rows == 0:
            log.error(f"Table '{S.table}' is empty or doesn't exist!")
            log.error("  You must index documents first (run without --query-only)")
            sys.exit(1)
        log.info(f"Using existing index with {existing_rows} rows")

        # Check if current parameters match the indexed data
        existing_config = check_index_configuration()
        current_signature = f"cs{S.chunk_size}_ov{S.chunk_overlap}_{S.embed_model_name.replace('/', '_')}"

        if existing_config:
            if existing_config.get("legacy") or not existing_config.get("has_metadata"):
                log.warning("  ⚠️  Querying LEGACY INDEX (no configuration metadata)")
                log.warning("  ⚠️  Cannot verify if chunk_size/overlap match query expectations")
            elif not existing_config.get("is_consistent"):
                log.warning("  ⚠️  MIXED INDEX: Table contains chunks from multiple configurations")
                log.warning("  ⚠️  Results may be inconsistent - consider rebuilding with RESET_TABLE=1")
            elif existing_config.get("index_signature") != current_signature:
                log.warning("  ⚠️  CONFIGURATION MISMATCH:")
                log.warning(f"    Query params: chunk_size={S.chunk_size}, overlap={S.chunk_overlap}")
                log.warning(f"    Index params: chunk_size={existing_config.get('chunk_size')}, overlap={existing_config.get('chunk_overlap')}")
                log.warning("  ⚠️  The CHUNK_SIZE you set doesn't affect retrieval - it only affects indexing!")
                log.warning("  ⚠️  You're querying chunks that were created with the index configuration above")
                log.warning(f"  💡 To query with chunk_size={S.chunk_size}, you need to re-index with RESET_TABLE=1")
            else:
                log.info(f"  ✓ Query configuration matches index: {existing_config.get('index_signature')}")

    # --- Build query engine ---
    log.info("=== LOADING LLM ===")
    llm = build_llm()

    # Choose retriever based on configuration
    log.info("=== CONFIGURING RETRIEVER ===")
    if S.hybrid_alpha < 1.0 or S.enable_filters or S.mmr_threshold > 0:
        log.info("Using HybridRetriever with advanced features:")
        if S.hybrid_alpha < 1.0:
            log.info(f"  ✓ Hybrid search: α={S.hybrid_alpha} (BM25 weight={1-S.hybrid_alpha:.1f})")
        if S.enable_filters:
            log.info(f"  ✓ Metadata filtering enabled (use 'participant:Name' or 'after:YYYY-MM-DD')")
        if S.mmr_threshold > 0:
            log.info(f"  ✓ MMR diversity enabled (λ={S.mmr_threshold})")

        retriever = HybridRetriever(
            vector_store,
            embed_model,
            similarity_top_k=S.top_k,
            alpha=S.hybrid_alpha,
            enable_metadata_filter=S.enable_filters,
            mmr_threshold=S.mmr_threshold
        )
    else:
        log.info("Using standard vector similarity search")
        retriever = VectorDBRetriever(vector_store, embed_model, similarity_top_k=S.top_k)

    query_engine = RetrieverQueryEngine.from_args(retriever, llm=llm)

    # --- Query mode selection ---
    if args.interactive:
        # Interactive REPL mode
        interactive_mode(query_engine, retriever)
    else:
        # Single query mode
        log.info("=== QUESTION ===")
        run_query(query_engine, S.question, retrieval_time=retriever.last_retrieval_time)


if __name__ == "__main__":
    main()
